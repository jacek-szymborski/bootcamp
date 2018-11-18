from pptx import Presentation
from csv_loader import srednia

prs = Presentation()

slide_layout = prs.slide_layouts[1]

slide = prs.slides.add_slide(slide_layout)
shapes = slide.shapes

title_shape = shapes.title

body_shape = shapes.placeholders[1]

title_shape.text = 'Jakis text'

tf = body_shape.text_frame
tf.text = 'Zawartosc text frame'

p = tf.add_paragraph()
p.text = "Kobiety - srednia wieku"
p.level = 1

p = tf.add_paragraph()
p.text = f"{srednia()}"
p.level = 2

prs.save("raport.pptx")

